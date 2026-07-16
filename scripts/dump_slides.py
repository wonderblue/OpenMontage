import sys
from pptx import Presentation

def dump_slides(pptx_path):
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides, 1):
        print(f"=== SLIDE {i} ===")
        # Get slide title if exists
        if slide.shapes.title:
            print(f"TITLE: {slide.shapes.title.text.strip()}")
        # Get all text blocks
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if slide.shapes.title and shape == slide.shapes.title:
                    continue
                print(shape.text.strip())
        print()

if __name__ == "__main__":
    if len(sys.argv) > 1:
        dump_slides(sys.argv[1])
    else:
        dump_slides("/Users/rkanduk/Documents/GitHub/OpenMontage/MM UNIT -1.pptx")
